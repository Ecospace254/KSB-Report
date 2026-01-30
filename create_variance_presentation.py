#!/usr/bin/env python3
"""
KSB Variance Analysis Presentation Generator
Creates a 5-slide PowerPoint presentation focusing on:
1. Sources of the 16.99% variance
2. Possible reasons for variances
3. Average yield, estimation, and difference to normal cane census
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
DARK_TEAL = RGBColor(0x00, 0x66, 0x66)
MEDIUM_TEAL = RGBColor(0x00, 0x88, 0x88)
LIGHT_TEAL = RGBColor(0x00, 0xAA, 0xAA)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
RED_ACCENT = RGBColor(0xC6, 0x28, 0x28)
ORANGE_ACCENT = RGBColor(0xE6, 0x5C, 0x00)
BLUE_ACCENT = RGBColor(0x15, 0x65, 0xC0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=14,
                       font_color=DARK_GRAY, bold_first=False, font_name="Calibri",
                       spacing=Pt(6), line_spacing=1.2):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = spacing
        p.line_spacing = line_spacing
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=DARK_TEAL, header_font_color=WHITE):
    """Add a formatted table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.text = data[0][col_idx]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = header_font_color
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for row_idx in range(1, rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
add_background(slide1, WHITE)

# Top accent bar
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Inches(0.15), DARK_TEAL)

# Left panel
add_shape(slide1, Inches(0), Inches(0.15), Inches(5.5), Inches(7.35), DARK_TEAL)

# Left panel content
add_textbox(slide1, Inches(0.5), Inches(1.0), Inches(4.5), Inches(1.0),
            "KENYA SUGAR BOARD", 16, WHITE, True, PP_ALIGN.LEFT)
add_textbox(slide1, Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.5),
            "& ECOSPACE SERVICES LTD.", 13, RGBColor(0xB0, 0xD0, 0xD0), False, PP_ALIGN.LEFT)

# Divider line on left
add_shape(slide1, Inches(0.5), Inches(2.2), Inches(2.5), Inches(0.04), WHITE)

add_textbox(slide1, Inches(0.5), Inches(2.5), Inches(4.5), Inches(1.2),
            "Drone-Assisted Sugarcane\nMapping Program", 20, WHITE, True, PP_ALIGN.LEFT)

add_textbox(slide1, Inches(0.5), Inches(3.9), Inches(4.5), Inches(0.8),
            "Western Catchment Region\nTechnical Report - January 2026", 14, RGBColor(0xB0, 0xD0, 0xD0),
            False, PP_ALIGN.LEFT)

# Bottom left info
add_textbox(slide1, Inches(0.5), Inches(6.2), Inches(4.5), Inches(0.5),
            "Ecospace Services Ltd.", 11, RGBColor(0x80, 0xB0, 0xB0), False, PP_ALIGN.LEFT)
add_textbox(slide1, Inches(0.5), Inches(6.5), Inches(4.5), Inches(0.5),
            "Providing Geospatial Solutions", 10, RGBColor(0x80, 0xB0, 0xB0), False, PP_ALIGN.LEFT)

# Right panel - main title
add_textbox(slide1, Inches(6.2), Inches(1.2), Inches(6.5), Inches(1.5),
            "VARIANCE ANALYSIS\nPRESENTATION", 36, DARK_TEAL, True, PP_ALIGN.LEFT)

# Subtitle
add_shape(slide1, Inches(6.2), Inches(3.0), Inches(3.0), Inches(0.04), DARK_TEAL)

add_textbox(slide1, Inches(6.2), Inches(3.3), Inches(6.5), Inches(1.0),
            "Understanding the 16.99% Census Overestimation:\nSources, Reasons & Yield Comparison",
            16, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# Key stats boxes
stats = [
    ("16.99%", "Census\nOverestimation"),
    ("92.3%", "Prediction\nAccuracy (R\u00B2)"),
    ("25,011 Ha", "Mature Cane\nIdentified"),
]
for i, (value, label) in enumerate(stats):
    x = Inches(6.2 + i * 2.2)
    y = Inches(4.8)
    box = add_shape(slide1, x, y, Inches(2.0), Inches(1.8), RGBColor(0xF5, 0xF5, 0xF5), MEDIUM_TEAL)
    add_textbox(slide1, x + Inches(0.1), y + Inches(0.2), Inches(1.8), Inches(0.8),
                value, 24, DARK_TEAL, True, PP_ALIGN.CENTER)
    add_textbox(slide1, x + Inches(0.1), y + Inches(1.0), Inches(1.8), Inches(0.6),
                label, 11, MEDIUM_GRAY, False, PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 2: SOURCES OF THE 16.99% VARIANCE
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)

# Header bar
add_shape(slide2, Inches(0), Inches(0), SLIDE_W, Inches(1.1), DARK_TEAL)
add_textbox(slide2, Inches(0.5), Inches(0.15), Inches(8), Inches(0.5),
            "SLIDE 1: SOURCES OF THE 16.99% VARIANCE", 26, WHITE, True)
add_textbox(slide2, Inches(0.5), Inches(0.6), Inches(10), Inches(0.4),
            "Drone-measured vs. Census-reported field areas across multiple sugar companies",
            13, RGBColor(0xB0, 0xD0, 0xD0), False)

# Intro text
add_textbox(slide2, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7),
            "Systematic comparison between drone-measured field areas and areas reported through conventional census "
            "enumeration revealed that traditional census methods overestimate mature cane area by an average of "
            "16.99% across the surveyed regions. The evidence was substantiated through field-by-field validation "
            "across multiple sugar companies and cooperatives:",
            12, DARK_GRAY, False)

# Table: Field-level variance evidence
table_data = [
    ["Sugar Company / Zone", "Reported Area (Ha)", "Measured Area (Ha)", "Overestimation (%)", "Variance Category"],
    ["Transmara - Keiyan Coop (Field 1)", "26.00", "8.82", "195%", "Extreme"],
    ["Transmara - Keiyan Coop (Field 2)", "13.00", "6.16", "111%", "Extreme"],
    ["Sukari Industries (Nyanza)", "2.50", "1.95", "28%", "High"],
    ["West Kenya (Kabras Zone)", "Varies", "Varies", "15-35%", "High"],
    ["Muhoroni (Oduo Zone)", "8.09", "7.75", "4%", "Low"],
    ["South Nyanza Sugar Co.", "Varied", "Varied", "~5-10%", "Moderate"],
    ["WEIGHTED AVERAGE", "--", "--", "16.99%", "Systematic"],
]

tbl = add_table(slide2, Inches(0.5), Inches(2.2), Inches(12.3), Inches(2.8),
                len(table_data), 5, table_data)

# Highlight the last row
table = tbl.table
for col_idx in range(5):
    cell = table.cell(len(table_data) - 1, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = RED_ACCENT

# Key insight boxes at bottom
insights = [
    ("Area Measurement Bias", "Farmers report total land holdings\nrather than the specific portion\nplanted to sugarcane"),
    ("Geographic Pattern", "Variance is highest in remote areas\n(Transmara: up to 195%) and lowest\nnear major mill operations (4%)"),
    ("Systematic Direction", "Overestimation is consistently\nupward -- never underestimation --\nindicating structural bias"),
    ("Operational Impact", "Mills plan crushing schedules on\nfigures inflated by ~1/6th, leading\nto persistent supply shortfalls"),
]
for i, (title, desc) in enumerate(insights):
    x = Inches(0.5 + i * 3.15)
    y = Inches(5.3)
    box = add_shape(slide2, x, y, Inches(2.95), Inches(1.9), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_textbox(slide2, x + Inches(0.15), y + Inches(0.1), Inches(2.65), Inches(0.4),
                title, 13, DARK_GREEN, True, PP_ALIGN.LEFT)
    add_shape(slide2, x + Inches(0.15), y + Inches(0.5), Inches(1.5), Inches(0.03), GREEN)
    add_textbox(slide2, x + Inches(0.15), y + Inches(0.6), Inches(2.65), Inches(1.2),
                desc, 10, DARK_GRAY, False, PP_ALIGN.LEFT)


# ============================================================================
# SLIDE 3: POSSIBLE REASONS FOR THE VARIANCES
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)

# Header
add_shape(slide3, Inches(0), Inches(0), SLIDE_W, Inches(1.1), DARK_TEAL)
add_textbox(slide3, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
            "SLIDE 2: POSSIBLE REASONS FOR THE VARIANCES", 26, WHITE, True)
add_textbox(slide3, Inches(0.5), Inches(0.6), Inches(10), Inches(0.4),
            "Why traditional census methods systematically overestimate sugarcane area by 16.99%",
            13, RGBColor(0xB0, 0xD0, 0xD0), False)

# Six reason cards - 2 rows of 3
reasons = [
    ("1", "Farmer Self-Reporting Bias",
     "Farmers often report total land holdings\n"
     "rather than the specific portion planted to\n"
     "sugarcane. This includes homestead areas,\n"
     "paths, and non-cane sections within their\n"
     "registered parcels.",
     RED_ACCENT),
    ("2", "Fallow & Failed Areas Included",
     "Census figures include fallow areas, failed\n"
     "sections, and field margins in reported\n"
     "acreages. Areas where cane died from\n"
     "drought, disease, or poor germination are\n"
     "still counted as 'under cane'.",
     ORANGE_ACCENT),
    ("3", "Boundary Estimation Errors",
     "Irregularly shaped smallholder plots\n"
     "(avg. 0.48 Ha) make visual area estimation\n"
     "inherently inaccurate. Errors compound\n"
     "in non-rectangular fields common in\n"
     "Western Kenya's fragmented landscape.",
     BLUE_ACCENT),
    ("4", "Strategic Area Inflation",
     "Some farmers intentionally inflate reported\n"
     "areas anticipating area-based payments,\n"
     "input allocations, or credit access. Larger\n"
     "reported areas may translate to greater\n"
     "perceived entitlements from mills.",
     RED_ACCENT),
    ("5", "Enumerator Limitations",
     "Census enumerators can only visit a small\n"
     "fraction of the total planted area within\n"
     "practical time and budget limits, leading\n"
     "to reliance on farmer declarations rather\n"
     "than independent verification.",
     ORANGE_ACCENT),
    ("6", "Temporal & Seasonal Gaps",
     "Manual census activities are seasonal and\n"
     "fail to capture rapid changes caused by\n"
     "planting cycles, weather variability, and\n"
     "disease outbreaks. A single snapshot cannot\n"
     "represent dynamic field conditions.",
     BLUE_ACCENT),
]

for i, (num, title, desc, color) in enumerate(reasons):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 3.1)

    # Card background
    card = add_shape(slide3, x, y, Inches(3.9), Inches(2.8), RGBColor(0xFA, 0xFA, 0xFA), RGBColor(0xE0, 0xE0, 0xE0))

    # Number circle
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide3, x + Inches(0.7), y + Inches(0.15), Inches(3.0), Inches(0.45),
                title, 14, color, True, PP_ALIGN.LEFT)

    # Divider
    add_shape(slide3, x + Inches(0.15), y + Inches(0.65), Inches(3.6), Inches(0.02), color)

    # Description
    add_textbox(slide3, x + Inches(0.15), y + Inches(0.75), Inches(3.6), Inches(1.9),
                desc, 11, DARK_GRAY, False, PP_ALIGN.LEFT)


# ============================================================================
# SLIDE 4: AVERAGE YIELD, ESTIMATION VS. NORMAL CANE CENSUS
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)

# Header
add_shape(slide4, Inches(0), Inches(0), SLIDE_W, Inches(1.1), DARK_TEAL)
add_textbox(slide4, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
            "SLIDE 3: AVERAGE YIELD - OUR ESTIMATION vs. NORMAL CANE CENSUS", 24, WHITE, True)
add_textbox(slide4, Inches(0.5), Inches(0.6), Inches(10), Inches(0.4),
            "Random Forest ML model (R\u00B2 = 0.923) compared against traditional census methodology",
            13, RGBColor(0xB0, 0xD0, 0xD0), False)

# LEFT COLUMN: Model Performance
add_shape(slide4, Inches(0.3), Inches(1.3), Inches(6.2), Inches(0.5), DARK_GREEN)
add_textbox(slide4, Inches(0.5), Inches(1.33), Inches(5.8), Inches(0.45),
            "DRONE/SATELLITE ESTIMATION (Our Model)", 15, WHITE, True, PP_ALIGN.LEFT)

metrics = [
    ("R\u00B2 (Coefficient of Determination)", "0.923", "92.3% of yield variation explained"),
    ("Mean Absolute Error (MAE)", "1.547 t/ha", "Predictions within 2-3% of actuals"),
    ("Mean Abs. Percentage Error (MAPE)", "6.5%", "Average prediction deviation"),
    ("Root Mean Square Error (RMSE)", "11.211 t/ha", "Larger outlier errors penalized"),
    ("Statistical Significance", "p < 0.001", "Highly significant results"),
    ("Classification Accuracy", "89.7%", "Sugarcane vs. other crops"),
    ("Sugarcane Producer's Accuracy", "94.2%", "Correctly identified cane fields"),
]

for i, (metric, value, note) in enumerate(metrics):
    y = Inches(1.95 + i * 0.55)
    bg_color = RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else WHITE
    add_shape(slide4, Inches(0.3), y, Inches(6.2), Inches(0.5), bg_color)
    add_textbox(slide4, Inches(0.4), y + Inches(0.05), Inches(2.8), Inches(0.4),
                metric, 10, DARK_GRAY, False, PP_ALIGN.LEFT)
    add_textbox(slide4, Inches(3.3), y + Inches(0.05), Inches(1.2), Inches(0.4),
                value, 11, DARK_GREEN, True, PP_ALIGN.CENTER)
    add_textbox(slide4, Inches(4.5), y + Inches(0.05), Inches(2.0), Inches(0.4),
                note, 9, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# RIGHT COLUMN: Census vs Drone Comparison
add_shape(slide4, Inches(6.8), Inches(1.3), Inches(6.2), Inches(0.5), RED_ACCENT)
add_textbox(slide4, Inches(7.0), Inches(1.33), Inches(5.8), Inches(0.45),
            "TRADITIONAL CANE CENSUS (Normal Method)", 15, WHITE, True, PP_ALIGN.LEFT)

census_items = [
    "Census area: 149,438 Ha reported across 10 counties",
    "Drone/satellite verified: ~124,046 Ha actual (after 16.99% correction)",
    "Difference: ~25,392 Ha of phantom cane area",
    "309,017 registered farmers at avg. 0.48 Ha plot size",
    "Enumerator-dependent: limited spatial coverage",
    "Seasonal snapshots: fail to capture rapid changes",
    "Farmer self-declared: no independent verification",
]

for i, item in enumerate(census_items):
    y = Inches(1.95 + i * 0.55)
    bg_color = RGBColor(0xFF, 0xEB, 0xEE) if i % 2 == 0 else WHITE
    add_shape(slide4, Inches(6.8), y, Inches(6.2), Inches(0.5), bg_color)
    add_textbox(slide4, Inches(6.9), y + Inches(0.05), Inches(6.0), Inches(0.4),
                item, 11, DARK_GRAY, False, PP_ALIGN.LEFT)

# Comparison section header
add_shape(slide4, Inches(0.3), Inches(5.95), SLIDE_W - Inches(0.6), Inches(0.45), MEDIUM_TEAL)
add_textbox(slide4, Inches(0.5), Inches(5.97), Inches(12), Inches(0.4),
            "INTERNATIONAL BENCHMARKING: How Our Model Compares", 14, WHITE, True, PP_ALIGN.LEFT)

# Comparison table
comp_data = [
    ["Metric", "Kenya (This Study)", "Australia", "Sri Lanka", "Global RF Mean", "Verdict"],
    ["R\u00B2 Accuracy", "0.923", "0.96", "0.91", "\u22480.96", "Competitive"],
    ["Method", "Random Forest + Drone", "Neural Networks", "Random Forest", "Various ML", "Cost-effective"],
    ["Data Source", "Sentinel-2 + DJI Mavic 3", "Multi-sensor", "Landsat-8/Sentinel", "Mixed", "Accessible"],
    ["Context", "Smallholder (0.48 Ha)", "Large commercial", "Mixed", "Various", "Harder challenge"],
]

add_table(slide4, Inches(0.3), Inches(6.5), SLIDE_W - Inches(0.6), Inches(0.85),
          len(comp_data), 6, comp_data, MEDIUM_TEAL)


# ============================================================================
# SLIDE 5: YIELD & SUPPLY-DEMAND IMPACT ANALYSIS
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)

# Header
add_shape(slide5, Inches(0), Inches(0), SLIDE_W, Inches(1.1), DARK_TEAL)
add_textbox(slide5, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
            "SLIDE 4: YIELD IMPACT & SUPPLY-DEMAND CONSEQUENCES", 24, WHITE, True)
add_textbox(slide5, Inches(0.5), Inches(0.6), Inches(10), Inches(0.4),
            "How the 16.99% variance cascades into the cane supply-demand crisis",
            13, RGBColor(0xB0, 0xD0, 0xD0), False)

# Mature Cane Distribution (left half)
add_shape(slide5, Inches(0.3), Inches(1.3), Inches(5.8), Inches(0.45), DARK_GREEN)
add_textbox(slide5, Inches(0.5), Inches(1.32), Inches(5.4), Inches(0.4),
            "MATURE CANE DISTRIBUTION (Drone-Verified)", 13, WHITE, True, PP_ALIGN.LEFT)

mature_data = [
    ["County", "Drone-Verified (Ha)", "% of Total"],
    ["Kakamega", "6,158.43", "24.6%"],
    ["Nandi", "4,272.19", "17.1%"],
    ["Uasin Gishu", "4,189.87", "16.7%"],
    ["Bungoma", "3,257.13", "13.0%"],
    ["Trans Nzoia", "3,165.97", "12.7%"],
    ["Busia", "2,548.07", "10.2%"],
    ["West Pokot", "1,204.76", "4.8%"],
    ["TOTAL", "25,011.17", "100%"],
]

tbl2 = add_table(slide5, Inches(0.3), Inches(1.85), Inches(5.8), Inches(3.0),
                 len(mature_data), 3, mature_data, DARK_GREEN)

# Highlight total row
t2 = tbl2.table
for col_idx in range(3):
    cell = t2.cell(len(mature_data) - 1, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xC8, 0xE6, 0xC9)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN

# Supply-Demand table (right half)
add_shape(slide5, Inches(6.5), Inches(1.3), Inches(6.5), Inches(0.45), RED_ACCENT)
add_textbox(slide5, Inches(6.7), Inches(1.32), Inches(6.1), Inches(0.4),
            "SUPPLY-DEMAND CRISIS (2025-2026 Crushing Season)", 13, WHITE, True, PP_ALIGN.LEFT)

sd_data = [
    ["Period", "Available (MT)", "Required (MT)", "Deficit (MT)", "Deficit %", "Mill Util."],
    ["Nov-Dec 2025", "1,069,936", "1,743,780", "673,844", "38.6%", "61.4%"],
    ["Jan-Mar 2026", "1,440,431", "2,615,670", "1,175,239", "44.9%", "55.1%"],
    ["Apr-Jun 2026", "1,309,151", "2,615,670", "1,306,519", "50.0%", "50.0%"],
    ["Jul-Oct 2026", "1,180,551", "3,487,560", "2,306,519", "66.1%", "33.9%"],
    ["ANNUAL", "5,000,069", "10,462,680", "5,461,621", "52.2%", "47.8%"],
]

tbl3 = add_table(slide5, Inches(6.5), Inches(1.85), Inches(6.5), Inches(2.3),
                 len(sd_data), 6, sd_data, RED_ACCENT)

# Highlight annual row
t3 = tbl3.table
for col_idx in range(6):
    cell = t3.cell(len(sd_data) - 1, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xCD, 0xD2)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = RED_ACCENT

# Bottom impact boxes
add_shape(slide5, Inches(6.5), Inches(4.3), Inches(6.5), Inches(0.35), ORANGE_ACCENT)
add_textbox(slide5, Inches(6.7), Inches(4.32), Inches(6.1), Inches(0.3),
            "ECONOMIC IMPACT OF THE VARIANCE", 12, WHITE, True, PP_ALIGN.LEFT)

impact_items = [
    "5.46 million tonne annual deficit = KShs 24.6 BILLION in foregone farmer income",
    "Mills operating at <50% capacity on average (47.8% utilization)",
    "The 16.99% overestimation means mills plan for ~25,392 Ha that don't exist",
    "Per-unit processing costs escalate as fixed costs spread across reduced output",
    "At KShs 4,500/tonne, every 1% of phantom area = ~KShs 1.45 billion misallocation",
]
add_bullet_textbox(slide5, Inches(6.5), Inches(4.7), Inches(6.5), Inches(2.5),
                   [f"\u2022  {item}" for item in impact_items],
                   11, DARK_GRAY, font_name="Calibri", spacing=Pt(4))

# Feature importance section (bottom left)
add_shape(slide5, Inches(0.3), Inches(5.0), Inches(5.8), Inches(0.35), BLUE_ACCENT)
add_textbox(slide5, Inches(0.5), Inches(5.02), Inches(5.4), Inches(0.3),
            "FEATURE IMPORTANCE: What Drives Yield Variance", 12, WHITE, True, PP_ALIGN.LEFT)

fi_data = [
    ["Variable", "Importance", "Implication"],
    ["Drone-measured Area", "33%", "Accurate boundaries are the #1 factor"],
    ["Historical Tonnes/Ha", "23%", "Past performance predicts future yield"],
    ["Yield Threshold (Crop Cycle)", "20%", "Ratoon age explains 1/5 of variation"],
    ["NDRE (Red Edge Index)", "7%", "Best spectral index for dense canopy"],
    ["GCI + NDVI + GNDVI", "15%", "Combined vegetation health indicators"],
    ["WDRVI + Crop Class", "2%", "Supplementary classification features"],
]

add_table(slide5, Inches(0.3), Inches(5.45), Inches(5.8), Inches(1.95),
          len(fi_data), 3, fi_data, BLUE_ACCENT)


# ============================================================================
# SLIDE 6 (5th content slide): RECOMMENDATIONS & WAY FORWARD
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, WHITE)

# Header
add_shape(slide6, Inches(0), Inches(0), SLIDE_W, Inches(1.1), DARK_TEAL)
add_textbox(slide6, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
            "SLIDE 5: RECOMMENDATIONS & WAY FORWARD", 26, WHITE, True)
add_textbox(slide6, Inches(0.5), Inches(0.6), Inches(10), Inches(0.4),
            "Addressing the variance gap: from census bias to precision agriculture",
            13, RGBColor(0xB0, 0xD0, 0xD0), False)

# Three columns of recommendations
rec_sections = [
    ("IMMEDIATE ACTIONS\n(0-6 Months)", DARK_GREEN, [
        "Replace census self-reporting with\ndrone-verified boundary mapping",
        "Apply the 16.99% correction factor\nto all current census-based planning",
        "Prioritize Transmara & West Kenya\nzones where variance exceeds 35%",
        "Integrate drone-measured areas into\nmill crushing schedule algorithms",
        "Complete Western Catchment farm\ndigitization (149,438 Ha target)",
        "Train 5 additional KSB officers in\nGIS and remote sensing methods",
    ]),
    ("MEDIUM-TERM STRATEGY\n(6-18 Months)", BLUE_ACCENT, [
        "Deploy regional drone mapping hubs\n(Western, Nyanza, Coastal)",
        "Acquire satellite imagery for all\ncatchment regions nationally",
        "Implement seasonal recalibration\nprotocol for the ML model",
        "Add soil type, weather data, and\npest/disease inputs to model",
        "Build real-time monitoring dashboard\nfor mill operations planning",
        "Reduce R3+ ratoon proportion from\n17% toward 10% benchmark",
    ]),
    ("LONG-TERM VISION\n(18-36 Months)", MEDIUM_TEAL, [
        "Achieve full national farm coverage\nacross all sugar-growing regions",
        "Implement continuous model retraining\nwith each harvest season's data",
        "Establish automated early-warning\nsystem for supply gap detection",
        "Target R\u00B2 improvement from 0.923\ntoward Australian benchmark (0.96)",
        "Reduce census overestimation from\n16.99% to < 5% nationally",
        "Position KSB as regional leader in\ndata-driven agricultural management",
    ]),
]

for col, (header, color, items) in enumerate(rec_sections):
    x = Inches(0.3 + col * 4.25)
    w = Inches(4.05)

    # Section header
    add_shape(slide6, x, Inches(1.3), w, Inches(0.7), color)
    add_textbox(slide6, x + Inches(0.1), Inches(1.33), w - Inches(0.2), Inches(0.65),
                header, 13, WHITE, True, PP_ALIGN.CENTER)

    # Items
    for i, item in enumerate(items):
        y = Inches(2.15 + i * 0.82)
        bg = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
        card = add_shape(slide6, x, y, w, Inches(0.72), bg, RGBColor(0xE0, 0xE0, 0xE0))

        # Bullet number
        bullet = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.08), y + Inches(0.15), Inches(0.35), Inches(0.35))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = color
        bullet.line.fill.background()
        tf = bullet.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide6, x + Inches(0.5), y + Inches(0.05), w - Inches(0.6), Inches(0.62),
                    item, 10, DARK_GRAY, False, PP_ALIGN.LEFT)

# Bottom summary bar
add_shape(slide6, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), DARK_TEAL)
add_textbox(slide6, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.35),
            "Kenya Sugar Board  |  Ecospace Services Ltd.  |  Drone-Assisted Sugarcane Mapping  |  "
            "Western Catchment Region  |  January 2026",
            10, WHITE, False, PP_ALIGN.CENTER)

# ── Save ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "KSB_Variance_Analysis_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
